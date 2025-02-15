from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Intelligent Document Management System"
subtitle.text = "AI-powered classification, search, and metadata extraction\nHackathon 2025\nTeam Name & Members"

# Slide 2: Problem Statement
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Problem Statement"
content.text = ("• Modern businesses struggle with large volumes of unstructured documents.\n"
                "• Manual document classification and search are time-consuming and inefficient.\n"
                "• Industries like legal and healthcare require intelligent solutions.\n\n"
                "Key Challenges:\n"
                "✓ Lack of automation in document classification.\n"
                "✓ Extracting key metadata without human intervention.\n"
                "✓ Poor searchability due to unstructured data.")

# Slide 3: Solution Overview
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Solution Overview"
content.text = ("• AI-powered document classification and tagging.\n"
                "• OCR-based metadata extraction.\n"
                "• Semantic understanding for document relationships.\n"
                "• Industry-specific document organization.\n"
                "• AI-driven search and collaboration tools.")

# Slide 4: Key Features
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Key Features"
content.text = ("1. Automated Classification & Tagging.\n"
                "2. Intelligent Content Extraction & Metadata Generation.\n"
                "3. Semantic Understanding & Relationship Discovery.\n"
                "4. Niche Document Organization.\n"
                "5. Innovative Document Management & Collaboration.")

# Slide 5: Tech Stack
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Tech Stack"
content.text = ("• Frontend: Flask\n"
                "• Backend: Flask, Python\n"
                "• Database: SQLite\n"
                "• AI/ML: TensorFlow, PyTorch, spaCy, Tesseract OCR\n"
                "• Search & Indexing: Elasticsearch\n"
                "• Collaboration: Flask-SocketIO\n"
                "• API Integration: DeepSeek API")

# Slide 6: System Architecture
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "System Architecture"
content.text = ("1. User uploads a document → AI processes it.\n"
                "2. Automated classification & tagging.\n"
                "3. Metadata extraction (dates, amounts, names).\n"
                "4. Storage & indexing in SQLite + Elasticsearch.\n"
                "5. AI-powered search and collaboration tools.")

# Slide 7: Implementation & Demo
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Implementation & Demo"
content.text = ("• Live Demo (if possible).\n"
                "• Screenshots of key features:\n"
                "  ✓ Document upload and classification.\n"
                "  ✓ Metadata extraction.\n"
                "  ✓ AI-powered search in action.\n"
                "  ✓ Version control and collaboration.")

# Slide 8: Future Scope & Scalability
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Future Scope & Scalability"
content.text = ("• Cloud Deployment (AWS/GCP).\n"
                "• Fine-Tuned AI Models.\n"
                "• Multi-Language Support.\n"
                "• Blockchain for Document Integrity.\n"
                "• Mobile App Version.")

# Slide 9: Challenges & Learnings
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Challenges & Learnings"
content.text = ("• Challenges Faced:\n"
                "  ✓ AI model optimization.\n"
                "  ✓ Handling large-scale document search.\n"
                "  ✓ Implementing real-time collaboration.\n\n"
                "• Solutions:\n"
                "  ✓ Fine-tuned NLP models.\n"
                "  ✓ Used Elasticsearch for fast indexing.\n"
                "  ✓ Optimized Flask-SocketIO for collaboration.")

# Slide 10: Conclusion & Q&A
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion & Q&A"
content.text = ("• AI-driven document management improves efficiency.\n"
                "• Faster workflows, reduced manual effort, better collaboration.\n"
                "• Future potential: scaling with cloud, blockchain, and mobile solutions.\n\n"
                "Thank you! Open for questions.")

# Save the PowerPoint file
pptx_filename = "Hackathon_Presentation.pptx"
prs.save(pptx_filename)

pptx_filename
